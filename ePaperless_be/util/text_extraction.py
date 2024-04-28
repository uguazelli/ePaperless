import pytesseract
from PIL import Image
from pdf2image import convert_from_path
import textract
import docx2txt
import pptx

def extract_text_from_document(file_path):
    print("Text extraction starting...")
    text = ""

    # Extract text from PDF
    if file_path.lower().endswith('.pdf'):
        pages = convert_from_path(file_path)
        for page in pages:
            text += pytesseract.image_to_string(page)
    
    # Extract text from image formats
    elif file_path.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
        with Image.open(file_path) as image:
            # Convert the image to PNG format
            image = image.convert('RGB')
            text += pytesseract.image_to_string(image)
    
    # Extract text from Word documents
    elif file_path.lower().endswith(('.doc', '.docx')):
        text += textract.process(file_path).decode('utf-8')
    
    # Extract text from PowerPoint presentations
    elif file_path.lower().endswith(('.ppt', '.pptx')):
        prs = pptx.Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
    
    # Extract text from text files
    elif file_path.lower().endswith('.txt'):
        with open(file_path, 'r') as file:
            text += file.read()
    
    # Add more conditions for other document types as needed

    else:
        print("Unsupported file format")
    
    print("Text extraction finished")
    return text